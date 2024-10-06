import os
import threading
import time

import pypandoc
import pythoncom
import win32com
from django.shortcuts import render
from django.http import HttpResponse, JsonResponse
from django.urls import reverse
import comtypes.client
from django.views.static import serve
from pptx import Presentation
import logging

from django.http import HttpResponse

from django.conf import settings

import tempfile

from django.views.decorators.csrf import csrf_exempt
from docx2pdf import convert

import os
from django.http import JsonResponse
from django.shortcuts import render
from django.urls import reverse
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

from .forms import UploadFileForm

import os
from django.shortcuts import render
from django.conf import settings
import subprocess

def convert_file(request):
    # Create temp directory if it doesn't exist
    temp_dir = 'temp'
    if not os.path.exists(temp_dir):
        os.makedirs(temp_dir)

    if request.method == 'POST':
        form = UploadFileForm(request.POST, request.FILES)
        if form.is_valid():
            # Get the uploaded file
            document = request.FILES['file']
            file_path = os.path.join(temp_dir, document.name)

            # Save the uploaded file
            with open(file_path, 'wb+') as destination:
                for chunk in document.chunks():
                    destination.write(chunk)

            # Initialize COM
            pythoncom.CoInitialize()
            try:
                # Convert the file to PDF
                output_path = os.path.splitext(file_path)[0] + '.pdf'
                convert(file_path, output_path)

                # Check if conversion succeeded
                if os.path.exists(output_path):
                    # به جای برگرداندن فایل، URL آن را برمی‌گردانیم
                    file_url = request.build_absolute_uri(reverse('download_pdf', args=[os.path.basename(output_path)]))
                    return JsonResponse({'file_url': file_url})
            finally:
                # Uninitialize COM
                pythoncom.CoUninitialize()

    else:
        form = UploadFileForm()

    return render(request, 'file_converter/file_converter.html', {'form': form})


def download_pdf(request, filename):
    file_path = os.path.join('temp', filename)
    if os.path.exists(file_path):
        with open(file_path, 'rb') as pdf_file:
            response = HttpResponse(pdf_file.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename="{filename}"'
        return response
    return HttpResponse("File not found", status=404)


#
# def download_pdf(request, filename):
#     file_path = os.path.join(settings.MEDIA_ROOT, 'temp', filename)
#     if os.path.exists(file_path):
#         return FileResponse(open(file_path, 'rb'), as_attachment=True, filename=filename)
#     return HttpResponse("File not found", status=404)

# def ppt_to_pdf(request):
#     if request.method == 'POST' and request.FILES['ppt_file']:
#         ppt_file = request.FILES['ppt_file']
#
#         # ذخیره فایل PowerPoint در یک مکان موقت
#         with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_ppt:
#             for chunk in ppt_file.chunks():
#                 temp_ppt.write(chunk)
#
#         # مسیر فایل PDF خروجی
#         pdf_path = os.path.join(settings.MEDIA_ROOT, 'temp.pdf')
#
#         # مقداردهی اولیه COM
#         pythoncom.CoInitialize()
#
#         try:
#             # تبدیل PowerPoint به PDF با استفاده از Microsoft PowerPoint
#             powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
#             powerpoint.Visible = 1
#
#             ppt = powerpoint.Presentations.Open(temp_ppt.name)
#             ppt.SaveAs(pdf_path, 32)  # 32 is the PDF format code
#             ppt.Close()
#             powerpoint.Quit()
#         finally:
#             # آزادسازی COM
#             pythoncom.CoUninitialize()
#
#         # پاک کردن فایل موقت PowerPoint
#         os.unlink(temp_ppt.name)
#
#         # دانلود فایل PDF
#         with open(pdf_path, 'rb') as pdf_file:
#             response = HttpResponse(pdf_file.read(), content_type='application/pdf')
#             response['Content-Disposition'] = 'attachment; filename="converted.pdf"'
#
#         # پاک کردن فایل PDF
#         os.unlink(pdf_path)
#
#         return response
#
#     return render(request, 'image_converter/ppt_to_pdf.html')

# def ppt_to_pdf(request):
#     if request.method == 'POST' and request.FILES.get('ppt_file'):
#         ppt_file = request.FILES['ppt_file']
#
#         # ذخیره فایل PowerPoint در یک مکان موقت
#         with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_ppt:
#             for chunk in ppt_file.chunks():
#                 temp_ppt.write(chunk)
#
#         # مسیر فایل PDF خروجی
#         pdf_filename = f"converted_{int(time.time())}.pdf"
#         pdf_path = os.path.join(settings.MEDIA_ROOT, 'temp', pdf_filename)
#
#         # مقداردهی اولیه COM
#         pythoncom.CoInitialize()
#
#         try:
#             # تبدیل PowerPoint به PDF با استفاده از Microsoft PowerPoint
#             powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
#             powerpoint.Visible = 1
#
#             ppt = powerpoint.Presentations.Open(temp_ppt.name)
#             ppt.SaveAs(pdf_path, 32)  # 32 is the PDF format code
#             ppt.Close()
#             powerpoint.Quit()
#         finally:
#             # آزادسازی COM
#             pythoncom.CoUninitialize()
#
#         # پاک کردن فایل موقت PowerPoint
#         os.unlink(temp_ppt.name)
#
#         # ایجاد URL برای دانلود و پیش‌نمایش PDF
#         pdf_url = request.build_absolute_uri(reverse('download_pdf', args=[pdf_filename]))
#
#         return JsonResponse({
#             'pdf_url': pdf_url,
#             'filename': pdf_filename
#         })
#
#     return render(request, 'image_converter/ppt_to_pdf.html')

# def ppt_to_pdf(request):
#     conversion_error = None
#
#     def convert_in_thread(temp_ppt_name, pdf_path):
#         nonlocal conversion_error
#         try:
#             pythoncom.CoInitialize()
#             powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
#             ppt = powerpoint.Presentations.Open(temp_ppt_name)
#             ppt.SaveAs(pdf_path, 32)  # 32 is the PDF format code
#             ppt.Close()
#             powerpoint.Quit()
#         except Exception as e:
#             print(f"Conversion error: {e}")  # Print error to console for debugging
#             conversion_error = str(e)  # Store the error message
#         finally:
#             pythoncom.CoUninitialize()
#
#     if request.method == 'POST' and request.FILES.get('ppt_file'):
#         ppt_file = request.FILES['ppt_file']
#
#         with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_ppt:
#             for chunk in ppt_file.chunks():
#                 temp_ppt.write(chunk)
#             temp_ppt_name = temp_ppt.name  # Store the temporary file name
#
#         pdf_filename = f"converted_{int(time.time())}.pdf"
#         pdf_dir = os.path.join(settings.MEDIA_ROOT, 'temp')  # Ensure 'temp' directory exists in your MEDIA_ROOT
#         pdf_path = os.path.join(pdf_dir, pdf_filename)
#
#         if not os.path.exists(pdf_dir):
#             os.makedirs(pdf_dir)
#
#         conversion_thread = threading.Thread(target=convert_in_thread, args=(temp_ppt_name, pdf_path))
#         conversion_thread.start()
#         conversion_thread.join()
#
#         try:  # Ensure temporary file is deleted even if an error occurs
#             os.unlink(temp_ppt_name)
#         except FileNotFoundError:
#             print("Temporary file not found for deletion")
#
#         if conversion_error:
#             return JsonResponse({'error': conversion_error}, status=500)
#
#         pdf_url = request.build_absolute_uri(
#             reverse('download_pdf', args=[pdf_filename]))  # Assuming you have a 'download_pdf' URL pattern
#
#         return JsonResponse({
#             'pdf_url': pdf_url,
#             'filename': pdf_filename
#         })
#
#     return render(request, 'file_converter/ppt_to_pdf.html')  # Replace with your template


# def ppt_to_pdf(request):
#     if request.method == 'POST':
#         form = UploadFileForm(request.POST, request.FILES)
#         if form.is_valid():
#             ppt_file = request.FILES['file']
#             ppt_file_path = os.path.join(settings.MEDIA_ROOT, ppt_file.name)
#             with open(ppt_file_path, 'wb+') as destination:
#                 for chunk in ppt_file.chunks():
#                     destination.write(chunk)
#
#             # مسیر فایل PDF خروجی
#             pdf_file_path = ppt_file_path.replace('.pptx', '.pdf')
#
#             # ایجاد PDF
#             c = canvas.Canvas(pdf_file_path, pagesize=letter)
#             presentation = Presentation(ppt_file_path)
#
#             # اضافه کردن اسلایدها به PDF
#             for slide in presentation.slides:
#                 c.drawString(100, 750, slide.shapes.title.text if slide.shapes.title else "No Title")
#                 c.showPage()  # ایجاد صفحه جدید برای هر اسلاید
#
#             c.save()
#
#             return render(request, 'file_converter/success.html', {'pdf_file': pdf_file_path})
#
#     else:
#         form = UploadFileForm()
#     return render(request, 'file_converter/ppt_to_pdf.html', {'form': form})
def ppt_to_pdf(request):
    if request.method == 'POST':
        form = UploadFileForm(request.POST, request.FILES)
        if form.is_valid():
            ppt_file = request.FILES['file']
            ppt_file_path = os.path.join(settings.MEDIA_ROOT, ppt_file.name)
            with open(ppt_file_path, 'wb+') as destination:
                for chunk in ppt_file.chunks():
                    destination.write(chunk)

            # مسیر فایل PDF خروجی
            pdf_file_path = ppt_file_path.replace('.pptx', '.pdf')

            # ایجاد PDF
            c = canvas.Canvas(pdf_file_path, pagesize=letter)
            presentation = Presentation(ppt_file_path)

            # اضافه کردن اسلایدها به PDF
            for slide in presentation.slides:
                c.drawString(100, 750, slide.shapes.title.text if slide.shapes.title else "No Title")
                c.showPage()  # ایجاد صفحه جدید برای هر اسلاید

            c.save()

            # ایجاد لینک دانلود
            pdf_file_url = os.path.basename(pdf_file_path)  # فقط نام فایل را بگیرید
            return render(request, 'file_converter/success.html', {'pdf_file_url': pdf_file_url})

    else:
        form = UploadFileForm()
    return render(request, 'file_converter/ppt_to_pdf.html', {'form': form})

def download_ppt_to_pdf_file(request, filename):
    file_path = os.path.join(settings.MEDIA_ROOT, filename)
    return serve(request, filename, document_root=settings.MEDIA_ROOT)